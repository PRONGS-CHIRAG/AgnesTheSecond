#!/usr/bin/env python3
"""Generate PITCH.pptx — AgnesTheSecond pitch deck.

Visual-first rebuild. Each slide has ONE strong moment, not a bullet dump:

  1. Title — 'Shows its work.'
  2. Problem — 3 broken options, punchy.
  3. Deterministic by design — pillar + hard-rule callout.
  4. The Prioritization Framework — rendered as visual bars, not a table.
  5. 124 / 125 — massive stat, nothing else on screen.
  6. Dot wall — 125 circles, 124 red + 1 green. The veto rate at a glance.
  7. Naive vs Agnes — side-by-side card contrast (sets up the live demo).
  8. Team + ask — 2x2 placeholder grid, final pitch line.

Speaker notes on every slide mirror the PITCH.md script verbatim.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Palette (mirror taim/insights/agnes.html) ──────────────────────────────
BG       = RGBColor(0x0A, 0x0F, 0x1A)
BG_ALT   = RGBColor(0x07, 0x0B, 0x14)
SURFACE  = RGBColor(0x11, 0x18, 0x27)
SURFACE2 = RGBColor(0x1A, 0x22, 0x36)
SURFACE3 = RGBColor(0x22, 0x2D, 0x42)
BORDER   = RGBColor(0x26, 0x33, 0x50)
TEXT     = RGBColor(0xE9, 0xF4, 0xFF)
TEXT2    = RGBColor(0x8B, 0xA3, 0xC7)
TEXT3    = RGBColor(0x5A, 0x74, 0x9A)
ACCENT   = RGBColor(0x6C, 0x63, 0xFF)
ACCENT2  = RGBColor(0x00, 0xD4, 0xAA)
WARN     = RGBColor(0xF5, 0x9E, 0x0B)
DANGER   = RGBColor(0xEF, 0x44, 0x44)
SUCCESS  = RGBColor(0x22, 0xC5, 0x5E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_HEAD = "Helvetica Neue"
FONT_BODY = "Helvetica Neue"
FONT_MONO = "Menlo"

# ── Helpers ────────────────────────────────────────────────────────────────


def set_bg(slide, prs, color=BG):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(
    slide,
    left, top, width, height,
    text,
    *,
    size=18,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT,
    font=FONT_BODY,
    line_spacing=None,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_rich_text(slide, left, top, width, height, runs, *, align=PP_ALIGN.LEFT):
    """runs = list[(text, size, bold, color, font)]."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for (t, size, bold, color, font) in runs:
        r = p.add_run()
        r.text = t
        r.font.name = font or FONT_BODY
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def rule(slide, left, top, width, color=ACCENT2, height=Inches(0.06)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def kicker(slide, text, *, color=ACCENT2):
    add_text(
        slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.35),
        text.upper(), size=11, bold=True, color=color, font=FONT_BODY,
    )


def title_big(slide, text):
    add_text(
        slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
        text, size=32, bold=True, color=TEXT, font=FONT_HEAD,
    )


def page_footer(slide, idx, total):
    add_text(
        slide, Inches(0.6), Inches(7.0), Inches(6), Inches(0.3),
        "AgnesTheSecond  ·  TUM.ai × Spherecast Makeathon 2026",
        size=9, color=TEXT3,
    )
    add_text(
        slide, Inches(11.0), Inches(7.0), Inches(2), Inches(0.3),
        f"{idx} / {total}",
        size=9, color=TEXT3, align=PP_ALIGN.RIGHT,
    )


def add_notes(slide, script):
    slide.notes_slide.notes_text_frame.text = script


def rounded(slide, l, t, w, h, *, fill=SURFACE, border=BORDER, border_w=1.0, radius=0.04):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.adjustments[0] = radius
    s.line.color.rgb = border
    s.line.width = Pt(border_w)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s


def chip(slide, l, t, text, *, fill=SURFACE2, border=BORDER, color=TEXT, bold=True, size=11, pad=0.15):
    w = Inches(pad * 2 + 0.11 * len(text) + 0.2)
    h = Inches(0.36)
    s = rounded(slide, l, t, w, h, fill=fill, border=border, radius=0.35)
    tf = s.text_frame
    tf.margin_left = Inches(pad)
    tf.margin_right = Inches(pad)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return s, w


# ── Slides ─────────────────────────────────────────────────────────────────


def slide_title(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)

    # Off-center stripe for visual character
    rule(slide, Inches(0.6), Inches(2.55), Inches(1.2), color=ACCENT2, height=Inches(0.09))

    add_text(
        slide, Inches(0.6), Inches(2.8), Inches(12), Inches(1.6),
        "AgnesTheSecond",
        size=84, bold=True, color=TEXT, font=FONT_HEAD, line_spacing=0.95,
    )
    add_text(
        slide, Inches(0.6), Inches(4.35), Inches(12), Inches(0.9),
        "Supply chain intelligence that shows its work.",
        size=26, bold=False, color=ACCENT2,
    )
    add_text(
        slide, Inches(0.6), Inches(5.15), Inches(12), Inches(0.5),
        "Maximum leverage.  Minimum concentration risk.  Deterministic by design.",
        size=15, color=TEXT2,
    )

    # Bottom row: team tagline on the left
    add_text(
        slide, Inches(0.6), Inches(6.85), Inches(7), Inches(0.3),
        "TUM.ai × Spherecast Makeathon 2026",
        size=11, bold=True, color=TEXT3,
    )
    # Page counter still matches other slides for consistency
    add_text(
        slide, Inches(11.0), Inches(6.85), Inches(2), Inches(0.3),
        f"{idx} / {total}",
        size=11, color=TEXT3, align=PP_ALIGN.RIGHT,
    )

    add_notes(
        slide,
        "Hi — we're the team behind Agnes The Second. We built a supply-chain "
        "intelligence platform for CPG companies that does two things most "
        "tools refuse to do: it balances cost against risk honestly, and every "
        "number we show you is a pure function of your data — not an LLM guessing.",
    )


def slide_problem(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    kicker(slide, "The problem", color=WARN)
    title_big(slide, "Procurement tools pick the wrong trade-off.")
    rule(slide, Inches(0.6), Inches(1.8), Inches(0.5))

    # 3 crossed-out options as big type, strikethrough feel
    options = [
        ('"Consolidate — save money."',       DANGER,  "legacy"),
        ('"Diversify — reduce risk."',        DANGER,  "legacy"),
        ('"Here\'s a GPT recommendation."',   DANGER,  "vibes"),
    ]
    for i, (text, color, tag) in enumerate(options):
        y = Inches(2.2 + i * 1.2)
        # × mark
        x_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), y + Inches(0.35), Inches(0.7), Inches(0.04)
        )
        x_shape.line.fill.background()
        x_shape.fill.solid()
        x_shape.fill.fore_color.rgb = color
        # Quote
        add_text(
            slide, Inches(1.55), y + Inches(0.1), Inches(10), Inches(0.7),
            text, size=28, bold=True, color=TEXT, font=FONT_HEAD,
        )
        # Tag
        add_text(
            slide, Inches(11.3), y + Inches(0.23), Inches(1.7), Inches(0.4),
            tag, size=11, bold=True, color=color, align=PP_ALIGN.RIGHT,
        )

    # Conclusion ribbon
    banner = rounded(
        slide, Inches(0.6), Inches(6.15), Inches(12.2), Inches(0.7),
        fill=SURFACE, border=WARN, border_w=1.25, radius=0.2,
    )
    tb = slide.shapes.add_textbox(
        Inches(0.85), Inches(6.25), Inches(11.8), Inches(0.5),
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "None of these survive a CFO conversation."
    r.font.name = FONT_HEAD
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WARN

    page_footer(slide, idx, total)
    add_notes(
        slide,
        "Procurement teams live inside a real tension. Consolidation gives you "
        "leverage. Diversification protects you when a supplier has a factory "
        "fire or a recall. Legacy tools tell you to pick one. The new wave of "
        "AI tools just hallucinate confident-sounding numbers with no audit "
        "trail. Neither approach survives a CFO conversation.",
    )


def slide_deterministic(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    kicker(slide, "Deterministic by design", color=ACCENT2)
    title_big(slide, "Every score is pure math.")
    add_text(
        slide, Inches(0.6), Inches(1.8), Inches(12), Inches(0.5),
        "Same inputs → same outputs.   Bit-for-bit.   Always.",
        size=18, color=ACCENT2,
    )

    # Left column: what's deterministic
    y0 = 2.7
    add_text(
        slide, Inches(0.6), Inches(y0), Inches(6), Inches(0.4),
        "PURE RULES", size=11, bold=True, color=TEXT3,
    )
    deterministic_items = [
        "Consolidation detection",
        "Risk assessment (5 risk types)",
        "5-dimension scoring framework",
        "Grade mapping + monopoly veto",
    ]
    for i, item in enumerate(deterministic_items):
        y = Inches(y0 + 0.45 + i * 0.55)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.12), Inches(0.18), Inches(0.18)
        )
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT2
        add_text(
            slide, Inches(0.95), y, Inches(5.5), Inches(0.5),
            item, size=17, color=TEXT,
        )

    # Right column: where LLMs are allowed
    add_text(
        slide, Inches(7.0), Inches(y0), Inches(6), Inches(0.4),
        "LLMs — SCOPED AND BOUNDED", size=11, bold=True, color=TEXT3,
    )
    llm_items = [
        ("SQL dispatch", "agent runs queries; data comes from DB, not model"),
        ("Evidence extraction", "grounded citations, schema-validated"),
        ("Prose polish", "summary text only — never scores or grades"),
    ]
    for i, (name, body) in enumerate(llm_items):
        y = Inches(y0 + 0.45 + i * 0.75)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(7.0), y + Inches(0.12), Inches(0.18), Inches(0.18)
        )
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        add_text(
            slide, Inches(7.35), y, Inches(5.8), Inches(0.4),
            name, size=17, bold=True, color=TEXT,
        )
        add_text(
            slide, Inches(7.35), y + Inches(0.4), Inches(5.8), Inches(0.4),
            body, size=13, color=TEXT2,
        )

    # Hard rule banner
    banner = rounded(
        slide, Inches(0.6), Inches(6.1), Inches(12.2), Inches(0.75),
        fill=SURFACE, border=ACCENT2, border_w=1.5, radius=0.2,
    )
    tb = slide.shapes.add_textbox(
        Inches(0.85), Inches(6.2), Inches(11.8), Inches(0.55),
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Hard rule in the code: "
    r.font.name = FONT_HEAD
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = ACCENT2
    r2 = p.add_run()
    r2.text = "the LLM can never change a score, a grade, or a supplier."
    r2.font.name = FONT_HEAD
    r2.font.size = Pt(16)
    r2.font.bold = True
    r2.font.color.rgb = TEXT

    page_footer(slide, idx, total)
    add_notes(
        slide,
        "Here's what separates us from the AI-wrapper tools. Every score Agnes "
        "shows you is a deterministic function of your data. Consolidation "
        "detection, risk scoring, the weighted framework — all pure Python. "
        "Same inputs give you the same outputs, byte-for-byte, on every run. "
        "168 tests enforce this. We use LLMs exactly where they help and "
        "nowhere they can hurt. The chat agent dispatches real SQL queries "
        "against your database — it doesn't invent answers, it routes "
        "questions to deterministic lookups. When we optionally polish a "
        "recommendation summary, the LLM can only rewrite the prose. It "
        "cannot change a grade, a score, or a supplier. That's a hard "
        "boundary in the code.",
    )


def slide_framework(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    kicker(slide, "The prioritization framework")
    title_big(slide, "5 dimensions.  One weighted score.  One grade.")
    rule(slide, Inches(0.6), Inches(1.8), Inches(0.5))

    dims = [
        ("Consolidation leverage",    0.35, ACCENT2, False),
        ("Evidence confidence",       0.25, ACCENT2, False),
        ("Compliance fit",            0.20, ACCENT2, False),
        ("Supplier diversification",  0.10, DANGER,  True),   # the key one
        ("Switching feasibility",     0.10, ACCENT2, False),
    ]
    max_w_in = 8.5  # scale width for 0.35 -> full bar
    bar_h = 0.45
    gap = 0.25
    y0 = 2.3
    for i, (label, weight, color, is_key) in enumerate(dims):
        y = Inches(y0 + i * (bar_h + gap))
        # Label
        add_text(
            slide, Inches(0.6), y + Inches(0.05), Inches(3.8), Inches(0.4),
            label, size=16, bold=is_key, color=TEXT if not is_key else DANGER,
        )
        # Key marker
        if is_key:
            add_text(
                slide, Inches(4.4), y + Inches(0.05), Inches(0.3), Inches(0.4),
                "⚙", size=16, bold=True, color=DANGER,
            )
        # Track
        track = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(4.7), y + Inches(0.12),
            Inches(max_w_in), Inches(bar_h * 0.45),
        )
        track.line.fill.background()
        track.fill.solid()
        track.fill.fore_color.rgb = SURFACE2
        # Fill
        fill_w = Inches(max_w_in * (weight / 0.35))
        fill = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(4.7), y + Inches(0.12),
            fill_w, Inches(bar_h * 0.45),
        )
        fill.line.fill.background()
        fill.fill.solid()
        fill.fill.fore_color.rgb = color
        # Weight label
        add_text(
            slide, Inches(13.333 - 0.6 - 0.7), y + Inches(0.05), Inches(0.7), Inches(0.4),
            f"{weight:.2f}", size=16, bold=True, color=color, font=FONT_MONO,
            align=PP_ALIGN.RIGHT,
        )

    # Veto callout (fits below the bars)
    veto_y = y0 + len(dims) * (bar_h + gap) + 0.25
    banner = rounded(
        slide, Inches(0.6), Inches(veto_y), Inches(12.2), Inches(1.05),
        fill=SURFACE, border=DANGER, border_w=1.5, radius=0.12,
    )
    tb = slide.shapes.add_textbox(
        Inches(0.95), Inches(veto_y + 0.15), Inches(11.6), Inches(0.85),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "⚙ Monopoly veto:  "
    r.font.name = FONT_HEAD
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = DANGER
    r2 = p.add_run()
    r2.text = (
        "if consolidation would leave ≤ 2 suppliers network-wide, "
        "the grade is automatically downgraded.  Not a heuristic — a hard rule."
    )
    r2.font.name = FONT_HEAD
    r2.font.size = Pt(16)
    r2.font.color.rgb = TEXT

    page_footer(slide, idx, total)
    add_notes(
        slide,
        "The scoring framework. Every recommendation is a weighted sum of five "
        "dimensions. Leverage, evidence, compliance, diversification, "
        "feasibility. You can see the weights. You can check the math. Here's "
        "the load-bearing rule: when consolidating would leave fewer than three "
        "suppliers for that ingredient across the entire network, the framework "
        "automatically downgrades the recommendation. Not a heuristic, not a "
        "prompt — a hard rule in the code.",
    )


def slide_big_stat(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, color=BG_ALT)
    kicker(slide, "The result, on real data", color=ACCENT2)

    # Massive 124 / 125
    add_text(
        slide, Inches(0), Inches(1.9), Inches(13.333), Inches(3.8),
        "124 / 125",
        size=260, bold=True, color=TEXT, align=PP_ALIGN.CENTER,
        font=FONT_HEAD, line_spacing=0.9,
    )
    # Subtitle
    add_text(
        slide, Inches(0), Inches(5.3), Inches(13.333), Inches(0.8),
        "consolidation opportunities trigger the concentration-risk veto.",
        size=22, color=ACCENT2, align=PP_ALIGN.CENTER,
    )
    # Annotation
    add_text(
        slide, Inches(0), Inches(6.1), Inches(13.333), Inches(0.5),
        "A naive tool would recommend consolidating all 125.  Agnes flags 124 for review.",
        size=14, color=TEXT2, align=PP_ALIGN.CENTER,
    )

    page_footer(slide, idx, total)
    add_notes(
        slide,
        "On our dataset — 61 companies, 357 ingredients, 40 suppliers, $1.52 "
        "billion in spend — 124 of 125 consolidation opportunities trigger "
        "the concentration-risk veto. A naive tool would tell you to "
        "consolidate all 125. Agnes tells you which one is actually safe.",
    )


def slide_dot_wall(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    kicker(slide, "125 consolidation opportunities")
    title_big(slide, "Each dot is one consolidation call.")
    rule(slide, Inches(0.6), Inches(1.8), Inches(0.5))

    # 25 cols x 5 rows = 125 dots
    cols = 25
    rows = 5
    total_dots = cols * rows
    area_left = Inches(0.8)
    area_right = Inches(12.5)
    area_w = area_right - area_left
    area_top = Inches(2.35)
    area_bottom = Inches(4.75)
    area_h = area_bottom - area_top

    dot_d = 0.32  # inches
    gap_x = (area_w - Inches(dot_d * cols)) / (cols - 1)
    gap_y = (area_h - Inches(dot_d * rows)) / (rows - 1)

    safe_index = 0  # single green dot — first position (Beta Carotene analog)
    for i in range(total_dots):
        r_idx = i // cols
        c_idx = i % cols
        x = area_left + c_idx * (Inches(dot_d) + gap_x)
        y = area_top + r_idx * (Inches(dot_d) + gap_y)
        color = SUCCESS if i == safe_index else DANGER
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, y, Inches(dot_d), Inches(dot_d)
        )
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = color

    # Legend
    legend_y = Inches(5.25)
    # Green legend
    g = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), legend_y + Inches(0.07), Inches(0.28), Inches(0.28)
    )
    g.line.fill.background()
    g.fill.solid()
    g.fill.fore_color.rgb = SUCCESS
    add_text(
        slide, Inches(1.2), legend_y, Inches(5.5), Inches(0.5),
        "1   safe_to_consolidate", size=16, bold=True, color=SUCCESS,
    )
    add_text(
        slide, Inches(1.2), legend_y + Inches(0.4), Inches(5.5), Inches(0.4),
        "Beta Carotene · 4 global suppliers · diversification green",
        size=12, color=TEXT2,
    )
    # Red legend
    r_s = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(7.2), legend_y + Inches(0.07), Inches(0.28), Inches(0.28)
    )
    r_s.line.fill.background()
    r_s.fill.solid()
    r_s.fill.fore_color.rgb = DANGER
    add_text(
        slide, Inches(7.6), legend_y, Inches(5.5), Inches(0.5),
        "124   review_required — keep a backup",
        size=16, bold=True, color=DANGER,
    )
    add_text(
        slide, Inches(7.6), legend_y + Inches(0.4), Inches(5.5), Inches(0.4),
        "≤ 2 global suppliers · full consolidation would create single-source risk",
        size=12, color=TEXT2,
    )

    # Tagline
    add_text(
        slide, Inches(0.6), Inches(6.45), Inches(12.2), Inches(0.5),
        "Maximum leverage.  Minimum concentration risk.",
        size=15, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER,
    )
    page_footer(slide, idx, total)
    add_notes(
        slide,
        "Visualised: every dot is one consolidation opportunity in the dataset. "
        "124 red, 1 green. A naive tool would recommend consolidating on all "
        "125. Agnes says: on 124 of these, keep a qualified backup supplier. "
        "That's the defensible position.",
    )


def draw_mock_card(slide, left, top, width, height, *, title_text, supplier,
                   grade_text, grade_color, score, dims, downgrade_banner=None,
                   accent_bar=ACCENT2):
    """Render a faux 'recommendation card' matching /agnes/ aesthetic."""
    card = rounded(slide, left, top, width, height, fill=SURFACE2, border=BORDER, radius=0.05)
    # Left accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height
    )
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_bar

    inner_left = left + Inches(0.32)
    inner_w = width - Inches(0.5)

    # Title line
    add_text(
        slide, inner_left, top + Inches(0.25), inner_w - Inches(2.4), Inches(0.4),
        title_text, size=14, bold=True, color=TEXT, font=FONT_HEAD,
    )
    # Grade chip (top right) + score pill
    gchip, gchip_w = chip(
        slide,
        left + width - Inches(2.8), top + Inches(0.22),
        grade_text,
        fill=SURFACE, border=grade_color, color=grade_color, size=10,
    )
    chip(
        slide,
        left + width - Inches(1.0), top + Inches(0.22),
        f"{score:.2f}",
        fill=SURFACE, border=ACCENT2, color=ACCENT2, size=10,
    )
    # Supplier line
    add_text(
        slide, inner_left, top + Inches(0.62), inner_w, Inches(0.35),
        f"→ {supplier}", size=11, color=ACCENT2,
    )

    # Downgrade banner if present
    next_y = top + Inches(1.0)
    if downgrade_banner:
        bn_h = Inches(0.75)
        bn = rounded(
            slide, inner_left, next_y, inner_w, bn_h,
            fill=SURFACE, border=WARN, border_w=1.0, radius=0.15,
        )
        add_text(
            slide, inner_left + Inches(0.15), next_y + Inches(0.07),
            inner_w - Inches(0.3), Inches(0.3),
            "⚙ Concentration-risk veto",
            size=10, bold=True, color=WARN,
        )
        add_text(
            slide, inner_left + Inches(0.15), next_y + Inches(0.32),
            inner_w - Inches(0.3), Inches(0.45),
            downgrade_banner,
            size=9, color=TEXT,
        )
        next_y = next_y + bn_h + Inches(0.15)

    # Dimension bars
    add_text(
        slide, inner_left, next_y, inner_w, Inches(0.3),
        "PRIORITIZATION FRAMEWORK", size=9, bold=True, color=ACCENT2,
    )
    next_y = next_y + Inches(0.32)
    label_w_in = 2.1
    bar_area_w_in = (inner_w.inches if hasattr(inner_w, 'inches') else (inner_w / 914400)) - label_w_in - 0.6
    for (key_label, val, col) in dims:
        # label
        add_text(
            slide, inner_left, next_y, Inches(label_w_in), Inches(0.25),
            key_label, size=9, color=TEXT2,
        )
        # track
        track_x = inner_left + Inches(label_w_in)
        track = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, track_x, next_y + Inches(0.08),
            Inches(bar_area_w_in), Inches(0.08),
        )
        track.line.fill.background()
        track.fill.solid()
        track.fill.fore_color.rgb = SURFACE3
        fill_w = Inches(bar_area_w_in * max(0.0, min(1.0, val)))
        if fill_w > Emu(0):
            f = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, track_x, next_y + Inches(0.08),
                fill_w, Inches(0.08),
            )
            f.line.fill.background()
            f.fill.solid()
            f.fill.fore_color.rgb = col
        # value
        add_text(
            slide, inner_left + Inches(label_w_in + bar_area_w_in + 0.05),
            next_y, Inches(0.5), Inches(0.25),
            f"{val:.2f}", size=9, color=TEXT, font=FONT_MONO,
        )
        next_y = next_y + Inches(0.3)


def slide_naive_vs_agnes(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    kicker(slide, "Same ingredient. Same data. Two different answers.")
    title_big(slide, "Naive tool   vs   Agnes.")
    rule(slide, Inches(0.6), Inches(1.8), Inches(0.5))

    # Column headers
    add_text(
        slide, Inches(0.6), Inches(2.1), Inches(6.0), Inches(0.4),
        "NAIVE TOOL   —   'consolidate!'",
        size=12, bold=True, color=DANGER,
    )
    add_text(
        slide, Inches(7.0), Inches(2.1), Inches(6.0), Inches(0.4),
        "AGNES   —   'review, keep a backup'",
        size=12, bold=True, color=ACCENT2,
    )

    # Naive card (left) — simplified + reckless
    naive_card = rounded(
        slide, Inches(0.6), Inches(2.55), Inches(6.0), Inches(3.8),
        fill=SURFACE2, border=BORDER, radius=0.05,
    )
    add_text(
        slide, Inches(0.85), Inches(2.75), Inches(5.5), Inches(0.5),
        "Microcrystalline Cellulose",
        size=16, bold=True, color=TEXT,
    )
    chip(
        slide, Inches(0.85), Inches(3.3), "HIGH PRIORITY",
        fill=SURFACE, border=DANGER, color=DANGER, size=10,
    )
    add_text(
        slide, Inches(0.85), Inches(3.9), Inches(5.5), Inches(0.5),
        "'13 companies buying from 2 suppliers.",
        size=13, color=TEXT2,
    )
    add_text(
        slide, Inches(0.85), Inches(4.2), Inches(5.5), Inches(0.5),
        "Consolidate to Colorcon. Save $X million.'",
        size=13, color=TEXT2,
    )
    add_text(
        slide, Inches(0.85), Inches(5.0), Inches(5.5), Inches(0.5),
        "— No tension matrix.",
        size=12, color=TEXT3,
    )
    add_text(
        slide, Inches(0.85), Inches(5.3), Inches(5.5), Inches(0.5),
        "— No diversification check.",
        size=12, color=TEXT3,
    )
    add_text(
        slide, Inches(0.85), Inches(5.6), Inches(5.5), Inches(0.5),
        "— No concentration-risk veto.",
        size=12, color=TEXT3,
    )
    add_text(
        slide, Inches(0.85), Inches(5.95), Inches(5.5), Inches(0.3),
        "— Confidence: \"trust me\"",
        size=12, bold=True, color=DANGER,
    )

    # Agnes card (right) — full detail
    draw_mock_card(
        slide,
        left=Inches(7.0), top=Inches(2.55), width=Inches(5.8), height=Inches(3.8),
        title_text="Microcrystalline Cellulose",
        supplier="Colorcon  (primary) + 1 qualified backup",
        grade_text="review_required", grade_color=WARN,
        score=0.82,
        downgrade_banner=(
            "Network has 2 global suppliers — full consolidation would create "
            "single-source risk. Keep a qualified backup."
        ),
        dims=[
            ("Consolidation",   0.70, ACCENT2),
            ("Evidence",        1.00, ACCENT2),
            ("Compliance",      1.00, ACCENT2),
            ("Diversification", 0.25, DANGER),
            ("Switching",       0.99, ACCENT2),
        ],
        accent_bar=WARN,
    )

    # Footer takeaway
    add_text(
        slide, Inches(0.6), Inches(6.55), Inches(12.2), Inches(0.5),
        "Same engine, different answer — reached by math you can defend to a CFO.",
        size=14, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER,
    )

    page_footer(slide, idx, total)
    add_notes(
        slide,
        "This slide is the setup for the live demo. Every other tool in this "
        "space looks like the card on the left: big number, confident "
        "recommendation, no audit trail. Agnes shows you the tension matrix, "
        "the grade, the veto reason, and the exact math. Same ingredient. "
        "Same database. Two completely different calls. After this slide I'll "
        "switch to the running app and walk through it live.",
    )


def slide_team_ask(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs)
    kicker(slide, "The team")
    title_big(slide, "Built by four people in a weekend.")
    rule(slide, Inches(0.6), Inches(1.8), Inches(0.5))

    # 2x2 team grid
    team = [
        ("[ TEAM MEMBER 1 ]", "Role · e.g. Backend / Engine",     "one-line tagline"),
        ("[ TEAM MEMBER 2 ]", "Role · e.g. Frontend / UX",        "one-line tagline"),
        ("[ TEAM MEMBER 3 ]", "Role · e.g. Data / Analytics",     "one-line tagline"),
        ("[ TEAM MEMBER 4 ]", "Role · e.g. Product / Pitch",      "one-line tagline"),
    ]
    card_w = Inches(6.0)
    card_h = Inches(1.55)
    for i, (name, role, tagline) in enumerate(team):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.3)
        y = Inches(2.2 + row * 1.75)
        card = rounded(slide, x, y, card_w, card_h, fill=SURFACE, border=BORDER, radius=0.06)
        # Avatar circle
        av = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.25), y + Inches(0.27), Inches(0.95), Inches(0.95)
        )
        av.line.fill.background()
        av.fill.solid()
        av.fill.fore_color.rgb = SURFACE2
        add_text(
            slide, x + Inches(0.25), y + Inches(0.5), Inches(0.95), Inches(0.4),
            name.split("]")[0].replace("[", "").strip()[:2] or "··",
            size=22, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER,
        )
        # Text
        add_text(
            slide, x + Inches(1.4), y + Inches(0.27), card_w - Inches(1.6), Inches(0.45),
            name, size=17, bold=True, color=TEXT,
        )
        add_text(
            slide, x + Inches(1.4), y + Inches(0.7), card_w - Inches(1.6), Inches(0.4),
            role, size=12, color=ACCENT2,
        )
        add_text(
            slide, x + Inches(1.4), y + Inches(1.02), card_w - Inches(1.6), Inches(0.4),
            tagline, size=11, color=TEXT2,
        )

    # Ask row
    banner = rounded(
        slide, Inches(0.6), Inches(5.85), Inches(12.2), Inches(1.15),
        fill=SURFACE, border=ACCENT2, border_w=1.5, radius=0.1,
    )
    add_text(
        slide, Inches(0.85), Inches(6.0), Inches(11.8), Inches(0.5),
        "Ready to pilot today.",
        size=22, bold=True, color=TEXT, font=FONT_HEAD,
    )
    add_text(
        slide, Inches(0.85), Inches(6.5), Inches(11.8), Inches(0.5),
        "Maximum leverage.   Minimum concentration risk.   Deterministic by design.",
        size=13, color=ACCENT2,
    )

    page_footer(slide, idx, total)
    add_notes(
        slide,
        "We built this in a weekend. Four people. One brain. Happy to walk you "
        "through the code, the data, or the math — whichever you want to dig "
        "into first. Thank you.",
    )


# ── Main ───────────────────────────────────────────────────────────────────


def main(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_problem,
        slide_deterministic,
        slide_framework,
        slide_big_stat,
        slide_dot_wall,
        slide_naive_vs_agnes,
        slide_team_ask,
    ]
    total = len(builders)
    for i, build in enumerate(builders, start=1):
        build(prs, i, total)

    prs.save(out_path)
    print(f"wrote {out_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    root = Path(__file__).resolve().parents[1]
    main(root / "PITCH.pptx")
